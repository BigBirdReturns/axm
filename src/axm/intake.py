"""
AXM Intake Layer - Universal Source Routing (v0.5.1)

Routes sources to the appropriate processing path:
- STRUCTURED (schema exists) → adapter path → high confidence
- UNSTRUCTURED (text soup) → compiler path → variable confidence

The key insight: when schema exists, we don't need to extract.
The schema IS the extraction. Confidence = 1.0.
"""

from __future__ import annotations

import hashlib
import json
import re
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum, auto
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple, Union
from abc import ABC, abstractmethod

from .coords import Coord, Major
from .ir import Node, Relation, Provenance, SourceInfo
from .program import Program, ProgramBuilder


# =============================================================================
# SOURCE TYPES
# =============================================================================

class SourceType(Enum):
    """Detected source type determines processing path."""
    
    # STRUCTURED (adapter path) - schema defines meaning
    XBRL = auto()           # SEC filings, financial reports
    FHIR = auto()           # Healthcare records (HL7)
    SCHEMA_ORG = auto()     # JSON-LD structured web data
    OPENAPI = auto()        # API specifications
    ICAL = auto()           # Calendar data
    RSS = auto()            # Feed data
    
    # SEMI-STRUCTURED (hybrid) - may have schema
    JSON = auto()           # Generic JSON
    XML = auto()            # Generic XML
    CSV = auto()            # Tabular data
    
    # UNSTRUCTURED (compiler path) - extract from text
    TEXT = auto()           # Plain text
    MARKDOWN = auto()       # Formatted text
    HTML = auto()           # Web pages (without JSON-LD)
    
    # CONTAINERS (extract content, then route)
    PDF = auto()
    DOCX = auto()
    XLSX = auto()
    PPTX = auto()


class ProcessingPath(Enum):
    """Which pipeline handles this source."""
    ADAPTER = "adapter"     # Structured: schema → nodes directly
    COMPILER = "compiler"   # Unstructured: text → lexer → parser → nodes


# =============================================================================
# DETECTION
# =============================================================================

class Detector:
    """Detect source type from content or metadata."""
    
    # File extension mappings
    EXT_MAP = {
        '.xml': SourceType.XML,
        '.xbrl': SourceType.XBRL,
        '.json': SourceType.JSON,
        '.jsonld': SourceType.SCHEMA_ORG,
        '.yaml': SourceType.OPENAPI,
        '.yml': SourceType.OPENAPI,
        '.ics': SourceType.ICAL,
        '.ical': SourceType.ICAL,
        '.rss': SourceType.RSS,
        '.atom': SourceType.RSS,
        '.csv': SourceType.CSV,
        '.tsv': SourceType.CSV,
        '.txt': SourceType.TEXT,
        '.md': SourceType.MARKDOWN,
        '.html': SourceType.HTML,
        '.htm': SourceType.HTML,
        '.pdf': SourceType.PDF,
        '.docx': SourceType.DOCX,
        '.doc': SourceType.DOCX,
        '.xlsx': SourceType.XLSX,
        '.xls': SourceType.XLSX,
        '.pptx': SourceType.PPTX,
    }
    
    # XML namespace patterns for schema detection
    XML_SCHEMAS = {
        'xbrl': SourceType.XBRL,
        'xmlns:xbrli': SourceType.XBRL,
        'us-gaap': SourceType.XBRL,
        'hl7.org/fhir': SourceType.FHIR,
        'fhir': SourceType.FHIR,
    }
    
    @classmethod
    def detect(cls, source: Union[str, Path, bytes]) -> Tuple[SourceType, str]:
        """
        Detect source type and extract content.
        
        Returns:
            (SourceType, content_string)
        """
        # Handle file paths
        if isinstance(source, Path):
            return cls._detect_file(source)
        
        if isinstance(source, str):
            # Check if it looks like a file path (short, no newlines, has extension)
            if len(source) < 500 and '\n' not in source:
                path = Path(source)
                if path.suffix and path.exists() and path.is_file():
                    return cls._detect_file(path)
            # Otherwise treat as content
            return cls._detect_content(source)
        
        # Handle raw bytes
        if isinstance(source, bytes):
            content = source.decode('utf-8', errors='ignore')
            return cls._detect_content(content)
    
    @classmethod
    def _detect_file(cls, path: Path) -> Tuple[SourceType, str]:
        """Detect from file extension and content."""
        ext = path.suffix.lower()
        
        # Container formats need extraction
        if ext == '.pdf':
            content = cls._extract_pdf(path)
            return SourceType.PDF, content
        
        if ext in ('.doc', '.docx'):
            content = cls._extract_docx(path)
            return SourceType.DOCX, content
        
        if ext in ('.xls', '.xlsx'):
            content = cls._extract_xlsx(path)
            return SourceType.XLSX, content
        
        if ext == '.pptx':
            content = cls._extract_pptx(path)
            return SourceType.PPTX, content
        
        # Text-based formats - read and sniff
        try:
            content = path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            content = path.read_text(encoding='latin-1')
        
        # Extension gives hint, content confirms
        hint = cls.EXT_MAP.get(ext, SourceType.TEXT)
        
        # Refine XML detection
        if hint == SourceType.XML:
            refined = cls._refine_xml(content)
            return refined, content
        
        # Refine JSON detection  
        if hint == SourceType.JSON:
            refined = cls._refine_json(content)
            return refined, content
        
        return hint, content
    
    @classmethod
    def _detect_content(cls, content: str) -> Tuple[SourceType, str]:
        """Detect from content alone."""
        stripped = content.strip()
        
        # XML family
        if stripped.startswith('<?xml') or stripped.startswith('<'):
            return cls._refine_xml(content), content
        
        # JSON family
        if stripped.startswith('{') or stripped.startswith('['):
            return cls._refine_json(content), content
        
        # Calendar
        if 'BEGIN:VCALENDAR' in content[:500]:
            return SourceType.ICAL, content
        
        # RSS/Atom
        if '<rss' in content[:500] or '<feed' in content[:500]:
            return SourceType.RSS, content
        
        # Markdown indicators
        if re.search(r'^#{1,6}\s', content[:500], re.MULTILINE):
            return SourceType.MARKDOWN, content
        
        # Default to text
        return SourceType.TEXT, content
    
    @classmethod
    def _refine_xml(cls, content: str) -> SourceType:
        """Refine XML type by checking namespaces."""
        header = content[:2000].lower()
        
        for pattern, source_type in cls.XML_SCHEMAS.items():
            if pattern in header:
                return source_type
        
        return SourceType.XML
    
    @classmethod
    def _refine_json(cls, content: str) -> SourceType:
        """Refine JSON type by checking for schemas."""
        try:
            data = json.loads(content)
        except json.JSONDecodeError:
            return SourceType.JSON
        
        if not isinstance(data, dict):
            return SourceType.JSON
        
        # JSON-LD / Schema.org
        if '@context' in data or '@type' in data:
            return SourceType.SCHEMA_ORG
        
        # OpenAPI
        if 'openapi' in data or 'swagger' in data:
            return SourceType.OPENAPI
        
        # FHIR
        if data.get('resourceType') in ('Bundle', 'Patient', 'Observation', 'Condition'):
            return SourceType.FHIR
        
        return SourceType.JSON
    
    # -------------------------------------------------------------------------
    # Container extraction (stubs - implement with actual libraries)
    # -------------------------------------------------------------------------
    
    @classmethod
    def _extract_pdf(cls, path: Path) -> str:
        """Extract text from PDF."""
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                return '\n'.join(page.extract_text() or '' for page in pdf.pages)
        except ImportError:
            pass
        
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            return '\n'.join(page.extract_text() or '' for page in reader.pages)
        except ImportError:
            pass
        
        # Fallback: return marker
        return f"[PDF: {path.name} - install pdfplumber or PyPDF2 to extract]"
    
    @classmethod
    def _extract_docx(cls, path: Path) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            doc = Document(path)
            return '\n'.join(para.text for para in doc.paragraphs)
        except ImportError:
            return f"[DOCX: {path.name} - install python-docx to extract]"
    
    @classmethod
    def _extract_xlsx(cls, path: Path) -> str:
        """Extract text from Excel spreadsheet."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"# {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    line = '\t'.join(str(c) if c is not None else '' for c in row)
                    if line.strip():
                        lines.append(line)
            return '\n'.join(lines)
        except ImportError:
            return f"[XLSX: {path.name} - install openpyxl to extract]"
    
    @classmethod
    def _extract_pptx(cls, path: Path) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        lines.append(shape.text)
            return '\n'.join(lines)
        except ImportError:
            return f"[PPTX: {path.name} - install python-pptx to extract]"


# =============================================================================
# ROUTING
# =============================================================================

class Router:
    """Route sources to adapter or compiler path."""
    
    # Types that definitely have schemas
    ADAPTER_TYPES = {
        SourceType.XBRL,
        SourceType.FHIR,
        SourceType.SCHEMA_ORG,
        SourceType.OPENAPI,
        SourceType.ICAL,
        SourceType.RSS,
    }
    
    # Types that definitely need extraction
    COMPILER_TYPES = {
        SourceType.TEXT,
        SourceType.MARKDOWN,
        SourceType.HTML,
        SourceType.PDF,
        SourceType.DOCX,
        SourceType.PPTX,
    }
    
    @classmethod
    def route(cls, source_type: SourceType, content: str) -> ProcessingPath:
        """Determine processing path."""
        
        if source_type in cls.ADAPTER_TYPES:
            return ProcessingPath.ADAPTER
        
        if source_type in cls.COMPILER_TYPES:
            return ProcessingPath.COMPILER
        
        # Semi-structured: inspect content
        if source_type == SourceType.JSON:
            if cls._json_has_schema(content):
                return ProcessingPath.ADAPTER
        
        if source_type == SourceType.XML:
            if cls._xml_has_schema(content):
                return ProcessingPath.ADAPTER
        
        if source_type == SourceType.CSV:
            # CSV with clear headers could go either way
            # For now, treat as compiler (extract values)
            pass
        
        if source_type == SourceType.XLSX:
            # Excel is usually semi-structured
            # Route to compiler for text extraction
            pass
        
        return ProcessingPath.COMPILER
    
    @classmethod
    def _json_has_schema(cls, content: str) -> bool:
        """Check if JSON has detectable schema."""
        try:
            data = json.loads(content)
            if isinstance(data, dict):
                # Has explicit schema markers
                if '@context' in data or 'openapi' in data or 'swagger' in data:
                    return True
                # Consistent structure (heuristic)
                if 'items' in data and isinstance(data['items'], list):
                    items = data['items']
                    if len(items) >= 2:
                        keys_0 = set(items[0].keys()) if isinstance(items[0], dict) else set()
                        if keys_0 and all(
                            isinstance(item, dict) and set(item.keys()) == keys_0 
                            for item in items[:5]
                        ):
                            return True
        except (json.JSONDecodeError, KeyError, TypeError):
            pass
        return False
    
    @classmethod
    def _xml_has_schema(cls, content: str) -> bool:
        """Check if XML has known schema."""
        header = content[:2000].lower()
        return any(
            pattern in header 
            for pattern in Detector.XML_SCHEMAS.keys()
        )


# =============================================================================
# ADAPTER BASE
# =============================================================================

@dataclass
class ExtractedEntity:
    """Entity extracted from structured source."""
    major: int
    type_: int
    subtype: int
    label: str
    value: Any = None
    unit: Optional[str] = None
    confidence: float = 1.0
    metadata: Dict[str, Any] = field(default_factory=dict)


class BaseAdapter(ABC):
    """Base class for structured data adapters."""
    
    SOURCE_TYPE: SourceType = None
    
    @abstractmethod
    def parse(self, content: str) -> List[ExtractedEntity]:
        """Parse content into entities."""
        pass
    
    def to_nodes(
        self, 
        entities: List[ExtractedEntity], 
        builder: ProgramBuilder,
        prov_id: str,
    ) -> List[Node]:
        """Convert entities to IR nodes."""
        nodes = []
        for entity in entities:
            coord = builder.coords.next(
                entity.major,
                entity.type_,
                entity.subtype,
            )
            node = Node(
                coord=coord,
                label=entity.label,
                prov_id=prov_id,
                value=entity.value,
                unit=entity.unit,
                metadata=entity.metadata,
            )
            nodes.append(node)
        return nodes


# =============================================================================
# ADAPTER REGISTRY
# =============================================================================

class AdapterRegistry:
    """Registry of available adapters."""
    
    _adapters: Dict[SourceType, type] = {}
    
    @classmethod
    def register(cls, source_type: SourceType):
        """Decorator to register an adapter."""
        def decorator(adapter_class):
            cls._adapters[source_type] = adapter_class
            return adapter_class
        return decorator
    
    @classmethod
    def get(cls, source_type: SourceType) -> Optional[BaseAdapter]:
        """Get adapter for source type."""
        adapter_class = cls._adapters.get(source_type)
        if adapter_class:
            return adapter_class()
        return None
    
    @classmethod
    def has(cls, source_type: SourceType) -> bool:
        """Check if adapter exists."""
        return source_type in cls._adapters
    
    @classmethod
    def list(cls) -> List[SourceType]:
        """List registered adapters."""
        return list(cls._adapters.keys())


# =============================================================================
# BUILT-IN ADAPTERS
# =============================================================================

@AdapterRegistry.register(SourceType.XBRL)
class XBRLAdapter(BaseAdapter):
    """Adapter for XBRL financial reports."""
    
    SOURCE_TYPE = SourceType.XBRL
    
    # XBRL concept → coordinate mapping
    CONCEPT_COORDS = {
        # Assets (Major 7 Quantity, type 1 Financial, subtype varies)
        'assets': (7, 1, 1),
        'currentassets': (7, 1, 2),
        'cash': (7, 1, 3),
        'cashandcashequivalents': (7, 1, 3),
        'inventory': (7, 1, 4),
        'accountsreceivable': (7, 1, 5),
        
        # Liabilities
        'liabilities': (7, 1, 10),
        'currentliabilities': (7, 1, 11),
        'longtermdebt': (7, 1, 12),
        'accountspayable': (7, 1, 13),
        
        # Equity
        'stockholdersequity': (7, 1, 20),
        'retainedearnings': (7, 1, 21),
        'commonstock': (7, 1, 22),
        
        # Income statement
        'revenues': (7, 2, 1),
        'revenue': (7, 2, 1),
        'salesrevenue': (7, 2, 2),
        'costofrevenue': (7, 2, 10),
        'grossprofit': (7, 2, 11),
        'operatingexpenses': (7, 2, 12),
        'operatingincome': (7, 2, 13),
        'netincome': (7, 2, 20),
        'earningspershare': (7, 2, 21),
    }
    
    def parse(self, content: str) -> List[ExtractedEntity]:
        """Parse XBRL into entities."""
        import xml.etree.ElementTree as ET
        
        entities = []
        
        try:
            root = ET.fromstring(content)
        except ET.ParseError as e:
            return entities
        
        # Find all elements with contextRef (these are facts)
        for elem in root.iter():
            context_ref = elem.get('contextRef')
            if context_ref is None:
                continue
            
            # Get concept name from tag
            tag = elem.tag
            if '}' in tag:
                concept = tag.split('}')[-1]
            else:
                concept = tag
            
            # Skip structural elements
            if concept.lower() in ('context', 'unit', 'schemaref'):
                continue
            
            value = elem.text
            if value is None:
                continue
            
            # Parse numeric value
            try:
                numeric_value = float(value.replace(',', ''))
            except (ValueError, TypeError):
                numeric_value = value
            
            # Get coordinates from concept
            concept_lower = concept.lower().replace('_', '').replace('-', '')
            coords = self.CONCEPT_COORDS.get(concept_lower, (7, 99, 1))
            
            # Get unit
            unit_ref = elem.get('unitRef')
            unit = 'USD' if unit_ref and 'usd' in unit_ref.lower() else None
            
            entities.append(ExtractedEntity(
                major=coords[0],
                type_=coords[1],
                subtype=coords[2],
                label=self._format_label(concept),
                value=numeric_value,
                unit=unit,
                confidence=1.0,
                metadata={
                    'xbrl_concept': concept,
                    'context_ref': context_ref,
                },
            ))
        
        return entities
    
    def _format_label(self, concept: str) -> str:
        """Format concept name as readable label."""
        # Add spaces before capitals
        label = re.sub(r'([a-z])([A-Z])', r'\1 \2', concept)
        return label


@AdapterRegistry.register(SourceType.ICAL)
class ICalAdapter(BaseAdapter):
    """Adapter for iCalendar data."""
    
    SOURCE_TYPE = SourceType.ICAL
    
    def parse(self, content: str) -> List[ExtractedEntity]:
        """Parse iCal into entities."""
        entities = []
        
        # Simple line-based iCal parsing
        current_event = {}
        in_event = False
        
        for line in content.split('\n'):
            line = line.strip()
            
            if line == 'BEGIN:VEVENT':
                in_event = True
                current_event = {}
            elif line == 'END:VEVENT':
                if current_event.get('SUMMARY'):
                    entities.append(ExtractedEntity(
                        major=6,  # Time
                        type_=1,  # Event
                        subtype=1,
                        label=current_event.get('SUMMARY', 'Untitled'),
                        value=None,
                        confidence=1.0,
                        metadata={
                            'dtstart': current_event.get('DTSTART'),
                            'dtend': current_event.get('DTEND'),
                            'location': current_event.get('LOCATION'),
                            'description': current_event.get('DESCRIPTION', '')[:200],
                        },
                    ))
                in_event = False
            elif in_event and ':' in line:
                # Handle property;params:value format
                key_part, _, value = line.partition(':')
                key = key_part.split(';')[0]
                current_event[key] = value
        
        return entities


@AdapterRegistry.register(SourceType.RSS)
class RSSAdapter(BaseAdapter):
    """Adapter for RSS/Atom feeds."""
    
    SOURCE_TYPE = SourceType.RSS
    
    def parse(self, content: str) -> List[ExtractedEntity]:
        """Parse RSS/Atom into entities."""
        import xml.etree.ElementTree as ET
        
        entities = []
        
        try:
            root = ET.fromstring(content)
        except ET.ParseError:
            return entities
        
        # RSS format
        for item in root.iter('item'):
            title = self._get_text(item, 'title')
            if title:
                entities.append(ExtractedEntity(
                    major=1,  # Entity
                    type_=10, # News/Content
                    subtype=1,
                    label=title,
                    confidence=1.0,
                    metadata={
                        'link': self._get_text(item, 'link'),
                        'pubDate': self._get_text(item, 'pubDate'),
                        'description': self._get_text(item, 'description', '')[:300],
                    },
                ))
        
        # Atom format
        for entry in root.iter('{http://www.w3.org/2005/Atom}entry'):
            title_elem = entry.find('{http://www.w3.org/2005/Atom}title')
            title = title_elem.text if title_elem is not None else None
            if title:
                link_elem = entry.find('{http://www.w3.org/2005/Atom}link')
                updated_elem = entry.find('{http://www.w3.org/2005/Atom}updated')
                entities.append(ExtractedEntity(
                    major=1,
                    type_=10,
                    subtype=1,
                    label=title,
                    confidence=1.0,
                    metadata={
                        'link': link_elem.get('href') if link_elem is not None else None,
                        'updated': updated_elem.text if updated_elem is not None else None,
                    },
                ))
        
        return entities
    
    def _get_text(self, elem, tag: str, default: str = None) -> Optional[str]:
        """Get text from child element."""
        child = elem.find(tag)
        return child.text if child is not None else default


# =============================================================================
# UNIVERSAL COMPILE
# =============================================================================

def compile_universal(
    source: Union[str, Path, bytes],
    config: Optional["Config"] = None,
) -> Program:
    """
    Universal compilation that auto-routes to adapter or compiler.
    
    Args:
        source: File path, URL, or raw content
        config: Optional compiler config
        
    Returns:
        Compiled Program
        
    Example:
        # Works with any input
        program = compile_universal("report.pdf")
        program = compile_universal("data.xbrl")
        program = compile_universal(raw_text)
    """
    from .compiler import Compiler, Config
    
    # Detect source type
    source_type, content = Detector.detect(source)
    path = Router.route(source_type, content)
    
    # Build source info
    source_uri = str(source) if isinstance(source, (str, Path)) else f"{source_type.name.lower()}://"
    source_info = SourceInfo(
        uri=source_uri,
        hash_sha256=hashlib.sha256(content.encode()).hexdigest(),
    )
    
    if path == ProcessingPath.ADAPTER:
        return _compile_via_adapter(source_type, content, source_info)
    else:
        return _compile_via_compiler(content, source_info, config)


def _compile_via_adapter(
    source_type: SourceType,
    content: str,
    source_info: SourceInfo,
) -> Program:
    """Structured data path - high confidence, no LLM needed."""
    
    adapter = AdapterRegistry.get(source_type)
    if adapter is None:
        # Fallback to compiler if no adapter
        return _compile_via_compiler(content, source_info, None)
    
    # Parse structured data
    entities = adapter.parse(content)
    
    # Build program
    builder = ProgramBuilder(source_info)
    
    # Create provenance for adapter extraction
    prov = Provenance(
        prov_id="prov_adapter_001",
        chunk_id="adapter",
        extractor=f"adapter:{source_type.name.lower()}",
        timestamp=datetime.utcnow().isoformat() + "Z",
        tier=0,  # Structured = Tier 0 (highest confidence)
        confidence=1.0,
    )
    builder.add_provenance(prov)
    
    # Convert entities to nodes
    nodes = adapter.to_nodes(entities, builder, prov.prov_id)
    for node in nodes:
        builder.add_node(node, deduplicate=True)
    
    return builder.build()


def _compile_via_compiler(
    content: str,
    source_info: SourceInfo,
    config: Optional["Config"],
) -> Program:
    """Unstructured text path - variable confidence, uses extraction."""
    from .compiler import Compiler, Config
    
    compiler = Compiler(config or Config.no_llm())
    return compiler.compile(content, source_type="string")


# =============================================================================
# MERGE
# =============================================================================

def merge_programs(
    programs: List[Program],
    priority: str = "adapter",
) -> Program:
    """
    Merge multiple programs into one.
    
    Args:
        programs: List of programs to merge
        priority: "adapter" (high confidence wins) or "latest" (last wins)
        
    Returns:
        Merged program
    """
    if not programs:
        raise ValueError("No programs to merge")
    
    if len(programs) == 1:
        return programs[0]
    
    # Use first program's source as base
    base_source = programs[0].source
    builder = ProgramBuilder(base_source)
    
    # Track content hashes to detect conflicts
    seen_hashes: Dict[str, Tuple[Node, float]] = {}  # hash -> (node, confidence)
    
    for program in programs:
        # Add all provenance (with unique IDs)
        for prov in program.provenance.values():
            # Make prov_id unique across programs
            unique_prov_id = f"{prov.prov_id}_{id(program)}"
            unique_prov = Provenance(
                prov_id=unique_prov_id,
                chunk_id=prov.chunk_id,
                extractor=prov.extractor,
                timestamp=prov.timestamp,
                tier=prov.tier,
                confidence=prov.confidence,
                source_start=prov.source_start,
                source_end=prov.source_end,
                model=prov.model,
                prompt_hash=prov.prompt_hash,
            )
            try:
                builder.add_provenance(unique_prov)
            except ValueError:
                pass  # Skip if already exists
        
        # Add nodes with conflict resolution
        for node in program.nodes.values():
            content_hash = node.content_hash
            prov = program.provenance.get(node.prov_id)
            confidence = prov.confidence if prov else 0.5
            
            # Update prov_id to unique version
            unique_prov_id = f"{node.prov_id}_{id(program)}"
            
            if content_hash in seen_hashes:
                existing_node, existing_conf = seen_hashes[content_hash]
                
                if priority == "adapter":
                    # Higher confidence wins
                    if confidence > existing_conf:
                        seen_hashes[content_hash] = (node, confidence)
                else:  # latest
                    # Replace with new
                    seen_hashes[content_hash] = (node, confidence)
            else:
                seen_hashes[content_hash] = (node, confidence)
        
        # Add all relations (will reference original node IDs)
        for rel in program.relations:
            try:
                builder.add_relation(rel)
            except ValueError:
                pass  # Duplicate, skip
    
    # Add winning nodes with fresh coordinates
    for node, confidence in seen_hashes.values():
        # Create new node with fresh coordinate from builder
        coord = builder.coords.next(
            node.coord.major,
            node.coord.type_,
            node.coord.subtype,
        )
        
        # Find the unique prov_id
        unique_prov_id = None
        for program in programs:
            if node.prov_id in program.provenance:
                unique_prov_id = f"{node.prov_id}_{id(program)}"
                break
        
        new_node = Node(
            coord=coord,
            label=node.label,
            prov_id=unique_prov_id or node.prov_id,
            value=node.value,
            unit=node.unit,
            metadata=node.metadata,
        )
        builder.add_node(new_node, deduplicate=False)
    
    return builder.build()


# =============================================================================
# HELPERS
# =============================================================================

def detect_source(source: Union[str, Path, bytes]) -> Tuple[SourceType, ProcessingPath]:
    """Detect source type and recommended path."""
    source_type, content = Detector.detect(source)
    path = Router.route(source_type, content)
    return source_type, path


def list_adapters() -> List[str]:
    """List available adapters."""
    return [st.name for st in AdapterRegistry.list()]


def avg_confidence(program: Program) -> float:
    """Calculate average confidence across all nodes."""
    if not program.nodes:
        return 0.0
    
    total = 0.0
    count = 0
    
    for node in program.nodes.values():
        prov = program.provenance.get(node.prov_id)
        if prov:
            total += prov.confidence
            count += 1
    
    return total / count if count > 0 else 0.0
